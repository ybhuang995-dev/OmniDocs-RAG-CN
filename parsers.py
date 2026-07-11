"""
OmniDocs RAG — Document Parsers

Responsible for:
- Reading 40+ file formats and converting to text
- Heading-aware chunking with overlap
- Code-aware chunking (Python AST, JS/TS regex)
- Auto-categorization (YAML frontmatter → H1 → filename)
"""

import re
import json
from pathlib import Path
from typing import Optional


# ──────────────────────────────────────────────
# Supported Extensions
# ──────────────────────────────────────────────
SUPPORTED_EXTENSIONS = {
    # Native text formats
    ".md", ".txt", ".rst", ".text", ".log",
    # Code files (wrapped in markdown for chunking)
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".scss",
    ".java", ".go", ".rs", ".c", ".cpp", ".h", ".hpp",
    ".rb", ".php", ".swift", ".kt", ".lua", ".sh", ".bash",
    # Data/config formats
    ".json", ".yaml", ".yml", ".toml", ".xml", ".csv", ".ini", ".cfg",
    # Web formats
    ".html", ".htm",
    # Binary document formats (require optional packages)
    ".pdf", ".docx", ".xlsx", ".pptx", ".ipynb",
}

_CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".scss",
    ".java", ".go", ".rs", ".c", ".cpp", ".h", ".hpp",
    ".rb", ".php", ".swift", ".kt", ".lua", ".sh", ".bash",
    ".json", ".yaml", ".yml", ".toml", ".xml", ".ini", ".cfg",
}


# [中文化] 共享工具：检测文本是否包含中文字符
# 检查前 100 个字符中是否含有 CJK 统一表意文字 (U+4E00 ~ U+9FFF)
# 这是一个轻量级检测，避免对整个大文本做全量扫描
def _is_chinese_text(text: str) -> bool:
    """检测文本是否以中文为主（前100字符中含CJK字符则视为中文）。"""
    return any('一' <= c <= '鿿' for c in text[:100])


# ──────────────────────────────────────────────
# Heading-Aware Chunking with Overlap + Code-Aware
# ──────────────────────────────────────────────
def _split_into_sentences(text: str) -> list[str]:
    """Split text into sentences for overlap."""
    # [中文化] 原正则 (?<=[.?!。])\s+ 有两个问题：
    #   1. 字符类缺少中文标点（！？；）
    #   2. \s+ 要求标点后必须有空格，但中文标点后直接接下一句
    # 修正：补全中文标点 + \s+ 改为 \s*（空格变为可选）
    pieces = re.split(r"(?<=[.!?。！？；])\s*|\n{2,}", text)
    return [s.strip() for s in pieces if s.strip()]


def _extract_sections(text: str, filepath: str) -> list[dict]:
    """Split markdown by headings (## and ###) into semantic chunks
    with 2-sentence overlap between consecutive chunks."""

    # Strip YAML frontmatter (\A = absolute start of string, safer than ^)
    text = re.sub(r"\A---.*?---\s*", "", text, flags=re.DOTALL)

    parts = re.split(r"(^#{1,3}\s+.+$)", text, flags=re.MULTILINE)

    raw_sections = []
    current_heading = "Introduction"
    current_h1 = Path(filepath).stem

    for part in parts:
        part = part.strip()
        if not part:
            continue

        heading_match = re.match(r"^(#{1,3})\s+(.+)$", part)
        if heading_match:
            level = len(heading_match.group(1))
            title = heading_match.group(2).strip()
            if level == 1:
                current_h1 = title
            current_heading = title
            continue

        clean = re.sub(r"[|\-\s#*`>]", "", part)
        if len(clean) < 30:
            continue

        # [中文化] 原逻辑按空白分词数判断是否需切分：
        #   - 英文：按空格切分后统计词数，>700 词则切
        #   - 中文：没有空格，整段被当成 1 个"词"，永远不触发切分
        # 修正：检测到中文 → 按字符数（2000 字为界），否则保持原逻辑
        if _is_chinese_text(part):
            # 中文：按字符数切分，每块约 1800 字（bge-m3 8192 token 内安全）
            if len(part) > 2000:
                sub_chunks = []
                for i in range(0, len(part), 1800):
                    sub_chunks.append(part[i:i + 1800])
            else:
                sub_chunks = [part]
        else:
            # 英文及其他：保持原逻辑按词数切分
            words = part.split()
            if len(words) > 700:
                sub_chunks = []
                for i in range(0, len(words), 600):
                    sub_chunk = " ".join(words[i:i + 600])
                    sub_chunks.append(sub_chunk)
            else:
                sub_chunks = [part]

        for idx, chunk_text in enumerate(sub_chunks):
            raw_sections.append({
                "text": chunk_text,
                "heading": current_heading,
                "parent_heading": current_h1,
                "source": filepath,
                "filename": Path(filepath).name,
                "sub_index": idx,
                "word_count": len(chunk_text.split())
            })

    # [中文化] 原逻辑取最后 2 句做重叠，但中文句子短（2 句可能仅 20-30 字）
    # 修正：中文取约 150 字重叠，英文保持 2 句
    for i in range(1, len(raw_sections)):
        if raw_sections[i]["source"] == raw_sections[i-1]["source"]:
            if _is_chinese_text(raw_sections[i-1]["text"]):
                # 中文：取前一个块最后约 150 字符
                prev_text = raw_sections[i-1]["text"]
                overlap = prev_text[-150:] if len(prev_text) > 150 else prev_text
            else:
                # 英文及其他：保持取最后 2 句
                prev_sentences = _split_into_sentences(raw_sections[i-1]["text"])
                overlap = " ".join(prev_sentences[-2:]) if len(prev_sentences) >= 2 else ""
            if overlap:
                raw_sections[i]["text"] = f"[...] {overlap}\n\n{raw_sections[i]['text']}"
                raw_sections[i]["word_count"] = len(raw_sections[i]["text"].split())

    return raw_sections


def _chunk_python_code(text: str, filepath: str) -> list[dict]:
    """Split Python source by top-level classes and functions via AST."""
    import ast
    chunks = []
    try:
        tree = ast.parse(text)
        lines = text.split("\n")
        for node in ast.iter_child_nodes(tree):
            if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
                start = node.lineno - 1
                end = node.end_lineno
                chunk_text = "\n".join(lines[start:end])
                node_type = type(node).__name__.replace("Def", "").replace("Async", "async ")
                chunks.append({
                    "text": f"```python\n{chunk_text}\n```",
                    "heading": f"{node_type}: {node.name}",
                    "parent_heading": Path(filepath).stem,
                    "source": filepath,
                    "filename": Path(filepath).name,
                    "sub_index": 0,
                    "word_count": len(chunk_text.split()),
                })
    except SyntaxError:
        return []  # fallback to standard chunking
    return chunks


def _chunk_js_code(text: str, filepath: str) -> list[dict]:
    """Split JS/TS source by function/class declarations via regex."""
    patterns = [
        r"((?:export\s+)?(?:async\s+)?function\s+\w+[^}]*\{(?:[^{}]|\{[^{}]*\})*\})",
        r"((?:export\s+)?class\s+\w+[^}]*\{(?:[^{}]|\{(?:[^{}]|\{[^{}]*\})*\})*\})",
        r"((?:export\s+)?(?:const|let|var)\s+\w+\s*=\s*(?:async\s+)?\([^)]*\)\s*=>\s*\{(?:[^{}]|\{[^{}]*\})*\})",
    ]
    chunks = []
    seen_ranges: list[tuple[int, int]] = []

    for pattern in patterns:
        for match in re.finditer(pattern, text, re.DOTALL):
            start, end = match.start(), match.end()
            # Skip if overlaps with already-found chunk
            if any(s <= start < e for s, e in seen_ranges):
                continue
            seen_ranges.append((start, end))
            chunk_text = match.group(1)
            # Extract name
            name_match = re.search(r"(?:function|class|const|let|var)\s+(\w+)", chunk_text)
            name = name_match.group(1) if name_match else "anonymous"
            node_type = "class" if "class " in chunk_text[:20] else "function"
            lang = Path(filepath).suffix.lstrip(".")
            chunks.append({
                "text": f"```{lang}\n{chunk_text}\n```",
                "heading": f"{node_type}: {name}",
                "parent_heading": Path(filepath).stem,
                "source": filepath,
                "filename": Path(filepath).name,
                "sub_index": 0,
                "word_count": len(chunk_text.split()),
            })

    return chunks


def _extract_sections_smart(text: str, filepath: str) -> list[dict]:
    """
    Smart chunking by file type:
    - .py     → by classes/functions via AST
    - .js/.ts → by functions via regex
    - others  → by headings (standard behavior)
    """
    suffix = Path(filepath).suffix.lower()

    if suffix == ".py":
        result = _chunk_python_code(text, filepath)
        if result:
            return result
    elif suffix in (".js", ".ts", ".jsx", ".tsx"):
        result = _chunk_js_code(text, filepath)
        if result:
            return result

    return _extract_sections(text, filepath)


def _categorize_file(filepath: str, content: str) -> str:
    """
    Auto-categorize file — three-level priority:
    1. YAML Frontmatter `category:` key (explicit user override)
    2. First H1 heading `# Title` in the document (dynamic, zero-effort)
    3. Filename stem as last resort
    """
    # Priority 1: YAML Frontmatter
    frontmatter_match = re.match(r"^---\s*\n(.*?)\n---", content, flags=re.DOTALL)
    if frontmatter_match:
        for line in frontmatter_match.group(1).split('\n'):
            if line.strip().lower().startswith('category:'):
                return line.split(':', 1)[1].strip().lower()

    # Priority 2: First H1 heading — "# My Document Title" → "my document title"
    h1_match = re.search(r"^#\s+(.+)$", content, flags=re.MULTILINE)
    if h1_match:
        raw_title = h1_match.group(1).strip()
        # Normalize: lowercase, strip markdown emphasis, keep letters/numbers/spaces
        category = re.sub(r"[*_`]", "", raw_title).lower()
        category = re.sub(r"[^\w\s-]", "", category).strip()
        if category:
            return category

    # Priority 3: Filename stem
    return Path(filepath).stem.lower().replace("_", " ").replace("-", " ")


# ──────────────────────────────────────────────
# Multi-Format File Reader
# ──────────────────────────────────────────────
def _read_file_to_text(filepath: str) -> Optional[str]:
    """Read a file and convert to markdown-like text for indexing."""
    ext = Path(filepath).suffix.lower()
    try:
        raw = Path(filepath).read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return None

    if not raw.strip():
        return None

    # Markdown / plain text — use as-is
    if ext in (".md", ".txt", ".text", ".rst", ".log"):
        return raw

    # HTML — strip tags to plain text
    if ext in (".html", ".htm"):
        try:
            from html.parser import HTMLParser
            class _StripHTML(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.parts = []
                    self._skip = False
                def handle_starttag(self, tag, attrs):
                    if tag in ("script", "style"):
                        self._skip = True
                def handle_endtag(self, tag):
                    if tag in ("script", "style"):
                        self._skip = False
                def handle_data(self, data):
                    if not self._skip:
                        self.parts.append(data)
            parser = _StripHTML()
            parser.feed(raw)
            return "\n".join(parser.parts)
        except Exception:
            return re.sub(r"<[^>]+>", "", raw)  # fallback: naive strip

    # Code files — wrap in markdown code block with filename heading
    if ext in _CODE_EXTENSIONS:
        lang = ext.lstrip(".")
        name = Path(filepath).name
        return f"# {name}\n\n```{lang}\n{raw}\n```"

    # ── Binary document formats (optional packages) ──

    # PDF via pypdf
    if ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(filepath)
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"## Page {i + 1}\n\n{text}")
            return "\n\n".join(pages) if pages else None
        except ImportError:
            return None  # pypdf not installed
        except Exception:
            return None

    # Word (.docx) via python-docx
    if ext == ".docx":
        try:
            import docx
            doc = docx.Document(filepath)
            lines = []
            for para in doc.paragraphs:
                style = para.style.name
                if style.startswith("Heading 1"):
                    lines.append(f"# {para.text}")
                elif style.startswith("Heading 2"):
                    lines.append(f"## {para.text}")
                elif style.startswith("Heading 3"):
                    lines.append(f"### {para.text}")
                elif para.text.strip():
                    lines.append(para.text)
            # Extract tables as markdown
            for table in doc.tables:
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if rows:
                    header = "| " + " | ".join(rows[0]) + " |"
                    sep = "|" + "|".join(["---"] * len(rows[0])) + "|"
                    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
                    lines.append(f"{header}\n{sep}\n{body}")
            return "\n\n".join(lines) if lines else None
        except ImportError:
            return None  # python-docx not installed
        except Exception:
            return None

    # Excel (.xlsx) via openpyxl
    if ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, data_only=True)
            sections = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        rows.append([str(c) if c is not None else "" for c in row])
                if rows:
                    header = "| " + " | ".join(rows[0]) + " |"
                    sep = "|" + "|".join(["---"] * len(rows[0])) + "|"
                    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:50])
                    sections.append(f"## Sheet: {sheet_name}\n\n{header}\n{sep}\n{body}")
            return "\n\n".join(sections) if sections else None
        except ImportError:
            return None  # openpyxl not installed
        except Exception:
            return None

    # PowerPoint (.pptx) via python-pptx
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text.strip() for shape in slide.shapes
                         if hasattr(shape, "text") and shape.text.strip()]
                if texts:
                    slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
            return "\n\n".join(slides) if slides else None
        except ImportError:
            return None  # python-pptx not installed
        except Exception:
            return None

    # Jupyter Notebook (.ipynb)
    if ext == ".ipynb":
        try:
            nb = json.loads(Path(filepath).read_text(encoding="utf-8"))
            sections = []
            for cell in nb.get("cells", []):
                source = "".join(cell.get("source", []))
                if not source.strip():
                    continue
                if cell["cell_type"] == "markdown":
                    sections.append(source)
                elif cell["cell_type"] == "code":
                    sections.append(f"```python\n{source}\n```")
            return "\n\n".join(sections) if sections else None
        except Exception:
            return None

    return raw
